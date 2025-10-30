# SPDX-License-Identifier: AGPL-3.0-or-later
"""Utilities for reading various unstructured file formats."""

from __future__ import annotations

import base64
import json
import os
import re
import xml.etree.ElementTree as ET
from email import policy
from email.parser import BytesParser
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

import yaml

# ---- soft deps are optional; we import lazily inside funcs ----

_TEXT_EXTENSIONS = {
    ".csv",
    ".cfg",
    ".conf",
    ".ini",
    ".log",
    ".md",
    ".properties",
    ".rst",
    ".text",
    ".txt",
    ".yaml",
    ".yml",
    ".toml",
    ".jsonl",
    ".ndjson",
    ".html",
    ".htm",
    ".xml",
}

_JSONL_MIMES = {"application/x-ndjson", "application/jsonl", "application/ndjson"}

_MAX_BYTES = int(float(os.getenv("SR_ADAPTER_MAX_BYTES", "524288000")))  # 500MB上限(可変)
_PREVIEW_N = int(os.getenv("SR_ADAPTER_PREVIEW_BYTES", "32"))


# ------------------------- helpers -------------------------

def _read_text_best_effort(path: Path) -> Tuple[str, Dict[str, object]]:
    """Try utf-8 first, then a few common encodings; if `charset_normalizer` is
    available, use it. Always returns text+meta without raising."""
    meta: Dict[str, object] = {}
    try:
        txt = path.read_text(encoding="utf-8", errors="strict")
        meta["encoding"] = "utf-8"
        return txt, meta
    except Exception:
        pass

    # try BOM/UTF-16/Shift_JIS quickly
    for enc in ("utf-8-sig", "utf-16", "cp932", "shift_jis", "latin-1"):
        try:
            txt = path.read_text(encoding=enc, errors="strict")
            meta["encoding"] = enc
            return txt, meta
        except Exception:
            continue

    # charset-normalizer があれば推定
    try:
        from charset_normalizer import from_path  # type: ignore
        res = from_path(str(path)).best()
        if res is not None:
            txt = str(res)
            meta["encoding"] = res.encoding
            meta["encoding_confidence"] = getattr(res, "encoding_aliases", None)
            return txt, meta
    except Exception:
        pass

    # 最終手段: utf-8 ignore
    txt = path.read_text(encoding="utf-8", errors="ignore")
    meta["encoding"] = "utf-8/ignore"
    return txt, meta


def _normalize_newlines(text: str) -> str:
    return text.replace("\r\n", "\n").replace("\r", "\n")


def _size_guard(path: Path) -> None:
    sz = path.stat().st_size
    if sz > _MAX_BYTES:
        raise ValueError(f"File too large: {sz} bytes > limit={_MAX_BYTES} bytes")


def _binary_preview(blob: bytes) -> Dict[str, object]:
    return {
        "binary_preview_bytes": min(len(blob), _PREVIEW_N),
        "binary_preview_base64": base64.b64encode(blob[:_PREVIEW_N]).decode("ascii"),
        "sha256": _sha256(blob),
        "length_bytes": len(blob),
    }


def _sha256(b: bytes) -> str:
    import hashlib
    return hashlib.sha256(b).hexdigest()


def _rtf_to_text(blob: str) -> str:
    """Very small RTF to text converter without external deps."""

    out: List[str] = []
    i = 0
    length = len(blob)
    while i < length:
        ch = blob[i]
        if ch == "\\":
            i += 1
            if i >= length:
                break
            control = blob[i]
            if control == "'" and i + 2 < length:
                hex_part = blob[i + 1 : i + 3]
                try:
                    out.append(bytes.fromhex(hex_part).decode("latin-1"))
                except Exception:
                    pass
                i += 3
                continue
            while i < length and blob[i].isalpha():
                i += 1
            if i < length and blob[i] in "-0123456789":
                i += 1
                while i < length and blob[i].isdigit():
                    i += 1
            if i < length and blob[i] == " ":
                i += 1
            continue
        if ch in "{}":
            i += 1
            continue
        out.append(ch)
        i += 1
    text = "".join(out)
    return _normalize_newlines(text)


def _pretty_xml_text(root: ET.Element) -> str:
    try:
        from xml.dom import minidom

        pretty = minidom.parseString(ET.tostring(root, encoding="utf-8")).toprettyxml(indent="  ")
        return "\n".join(line for line in pretty.splitlines() if line.strip())
    except Exception:
        return _normalize_newlines(ET.tostring(root, encoding="unicode"))


def _prepare_image_for_ocr(image: "Image.Image") -> Tuple["Image.Image", "Image.Image", bool]:  # type: ignore[name-defined]
    """Return an orientation-corrected copy and a contrast-enhanced grayscale copy."""

    corrected = image
    orientation_changed = False
    try:
        from PIL import ImageOps  # type: ignore

        corrected_candidate = ImageOps.exif_transpose(image)
        if corrected_candidate is not None:
            corrected = corrected_candidate
            orientation_changed = corrected is not image
    except Exception:
        corrected = image

    ocr_ready = corrected
    try:
        if ocr_ready.mode not in {"RGB", "RGBA", "L"}:
            ocr_ready = ocr_ready.convert("RGB")
    except Exception:
        pass

    try:
        if ocr_ready.mode != "L":
            ocr_ready = ocr_ready.convert("L")
    except Exception:
        pass

    try:
        from PIL import ImageEnhance  # type: ignore

        enhancer = ImageEnhance.Contrast(ocr_ready)
        ocr_ready = enhancer.enhance(1.4)
    except Exception:
        pass

    return corrected, ocr_ready, orientation_changed


def _select_ocr_languages(pytesseract: Any) -> List[str]:
    """Pick a stable set of OCR languages based on what Tesseract exposes."""

    available: List[str] = []
    try:
        langs = pytesseract.get_languages(config="")  # type: ignore[attr-defined]
        if isinstance(langs, (list, tuple)):
            available = list(langs)
    except Exception:
        try:
            langs = pytesseract.get_available_languages()  # type: ignore[attr-defined]
            if isinstance(langs, (list, tuple, set)):
                available = list(langs)
        except Exception:
            available = []

    if not available:
        return []

    preferred = [
        "eng",
        "jpn",
        "kor",
        "chi_sim",
        "chi_tra",
        "fra",
        "deu",
        "spa",
        "ita",
        "por",
        "nld",
        "ara",
        "hin",
        "rus",
    ]
    seen: List[str] = []
    for code in preferred:
        if code in available and code not in seen:
            seen.append(code)

    if not seen:
        for code in available:
            if code not in seen:
                seen.append(code)
            if len(seen) >= 5:
                break

    return seen


# ------------------------- readers -------------------------

def _read_json(path: Path) -> Tuple[str, Dict[str, object]]:
    raw = path.read_text(encoding="utf-8", errors="ignore")
    meta: Dict[str, object] = {"json_valid": False}
    try:
        data = json.loads(raw)
        meta["json_valid"] = True
        meta["json_top_level_type"] = type(data).__name__
        if isinstance(data, dict):
            meta["json_top_level_keys"] = sorted(map(str, data.keys()))
        # pretty-print for downstream readability
        return json.dumps(data, ensure_ascii=False, indent=2), meta
    except json.JSONDecodeError as e:
        meta["json_error"] = str(e)
        return raw, meta


def _read_jsonl(path: Path) -> Tuple[str, Dict[str, object]]:
    meta: Dict[str, object] = {}
    lines = path.read_text(encoding="utf-8", errors="ignore").splitlines()
    valid = 0
    samples: List[Dict[str, object]] = []
    for i, ln in enumerate(lines[:50]):  # sample前半50行だけ解析
        try:
            obj = json.loads(ln)
            valid += 1
            if isinstance(obj, dict):
                samples.append({"line": i + 1, "keys": sorted(map(str, obj.keys()))})
        except Exception:
            pass
    meta.update({"jsonl_lines": len(lines), "jsonl_valid_count": valid, "jsonl_key_samples": samples})
    # そのまま返す（改変せず）
    return "\n".join(lines), meta


def _read_yaml(path: Path) -> Tuple[str, Dict[str, object]]:
    raw, meta = _read_text_best_effort(path)
    text = _normalize_newlines(raw)
    try:
        docs = list(yaml.safe_load_all(text))
        meta.update(
            {
                "yaml_documents": len(docs),
                "yaml_sample_types": [type(doc).__name__ for doc in docs[:5]],
            }
        )
        if docs and isinstance(docs[0], dict):
            meta["yaml_root_keys"] = sorted(map(str, docs[0].keys()))
    except Exception as exc:
        meta["yaml_error"] = str(exc)
    return text, meta


def _read_toml(path: Path) -> Tuple[str, Dict[str, object]]:
    raw = path.read_text(encoding="utf-8", errors="ignore")
    text = _normalize_newlines(raw)
    meta: Dict[str, object] = {"toml_valid": False}
    try:
        try:
            import tomllib  # type: ignore[attr-defined]
        except ModuleNotFoundError:  # pragma: no cover - Python <3.11 fallback
            import tomli as tomllib  # type: ignore
    except ModuleNotFoundError:
        meta["toml_error"] = "tomllib not available"
        return text, meta

    try:
        data = tomllib.loads(text)
        meta["toml_valid"] = True
        if isinstance(data, dict):
            meta["toml_root_keys"] = sorted(map(str, data.keys()))
            meta["toml_table_count"] = sum(1 for value in data.values() if isinstance(value, dict))
    except Exception as exc:
        meta["toml_error"] = str(exc)
    return text, meta


def _read_plain_text(path: Path) -> Tuple[str, Dict[str, object]]:
    txt, meta = _read_text_best_effort(path)
    return _normalize_newlines(txt), meta


def _read_markdown(path: Path) -> Tuple[str, Dict[str, object]]:
    txt, meta = _read_text_best_effort(path)
    txt = _normalize_newlines(txt)
    # front matter (--- YAML ---) を拾う
    if txt.startswith("---\n"):
        try:
            import yaml  # type: ignore
            end = txt.find("\n---", 4)
            if end != -1:
                fm = txt[4:end]
                meta["front_matter"] = yaml.safe_load(fm) or {}
                txt = txt[end + 4 :]  # YAMLを本文から除去
        except Exception:
            meta["front_matter"] = "unparsed"
    return txt, meta


def _read_html(path: Path) -> Tuple[str, Dict[str, object]]:
    txt, meta = _read_text_best_effort(path)
    try:
        from bs4 import BeautifulSoup  # type: ignore
        soup = BeautifulSoup(txt, "html.parser")
        title = soup.title.string.strip() if soup.title and soup.title.string else None
        for tag in soup(["script", "style", "noscript"]):
            tag.extract()
        text = " ".join(soup.get_text(separator=" ").split())
        meta.update({"html_title": title, "html_length": len(txt)})
        return text, meta
    except Exception:
        meta["html_parsed"] = False
        return txt, meta


def _read_xml(path: Path) -> Tuple[str, Dict[str, object]]:
    meta: Dict[str, object] = {"xml_valid": False}
    try:
        tree = ET.parse(str(path))
        root = tree.getroot()
        meta.update(
            {
                "xml_valid": True,
                "xml_root": root.tag,
                "xml_child_count": sum(1 for _ in root),
                "xml_element_count": sum(1 for _ in root.iter()),
            }
        )
        pretty = _pretty_xml_text(root)
        return pretty, meta
    except ET.ParseError as exc:
        meta["xml_error"] = str(exc)
    except Exception as exc:
        meta["xml_error"] = f"{type(exc).__name__}: {exc}"
    raw = path.read_text(encoding="utf-8", errors="ignore")
    return _normalize_newlines(raw), meta


def _read_pdf(path: Path) -> Tuple[str, Dict[str, object]]:
    try:
        from pypdf import PdfReader  # type: ignore
        reader = PdfReader(str(path))
        pages = min(len(reader.pages), int(os.getenv("SR_ADAPTER_PDF_MAX_PAGES", "400")))
        out: List[str] = []
        for i in range(pages):
            try:
                out.append(reader.pages[i].extract_text() or "")
            except Exception:
                out.append("")
        meta = {"pdf_pages": len(reader.pages), "pdf_pages_parsed": pages}
        return _normalize_newlines("\n".join(out)), meta
    except Exception as e:
        # バイナリプレビューにフォールバック
        blob = path.read_bytes()
        meta = _binary_preview(blob)
        meta["pdf_text_extraction"] = f"failed: {type(e).__name__}"
        return "", meta


def _read_docx(path: Path) -> Tuple[str, Dict[str, object]]:
    try:
        import docx  # type: ignore
        d = docx.Document(str(path))
        paras = [p.text for p in d.paragraphs]
        meta = {"docx_paragraphs": len(paras)}
        return _normalize_newlines("\n".join(paras)), meta
    except Exception as e:
        blob = path.read_bytes()
        meta = _binary_preview(blob)
        meta["docx_text_extraction"] = f"failed: {type(e).__name__}"
        return "", meta


def _read_pptx(path: Path) -> Tuple[str, Dict[str, object]]:
    try:
        from pptx import Presentation  # type: ignore
        prs = Presentation(str(path))
        texts: List[str] = []
        for s in prs.slides:
            for shp in s.shapes:
                try:
                    if hasattr(shp, "text"):
                        texts.append(shp.text)
                except Exception:
                    pass
        meta = {"pptx_slides": len(prs.slides)}
        return _normalize_newlines("\n".join(texts)), meta
    except Exception as e:
        meta = _binary_preview(path.read_bytes())
        meta["pptx_text_extraction"] = f"failed: {type(e).__name__}"
        return "", meta


def _read_xlsx(path: Path) -> Tuple[str, Dict[str, object]]:
    try:
        import openpyxl  # type: ignore
        wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
        texts: List[str] = []
        total_rows = 0
        for ws in wb.worksheets:
            for r, row in enumerate(ws.iter_rows(values_only=True)):
                # 行数が極端に多いとき用の上限
                if r > int(os.getenv("SR_ADAPTER_XLSX_MAX_ROWS", "10000")):
                    break
                total_rows += 1
                vals = ["" if v is None else str(v) for v in row]
                texts.append("\t".join(vals))
        meta = {"xlsx_sheets": len(wb.worksheets), "xlsx_rows_read": total_rows}
        return _normalize_newlines("\n".join(texts)), meta
    except Exception as e:
        meta = _binary_preview(path.read_bytes())
        meta["xlsx_text_extraction"] = f"failed: {type(e).__name__}"
        return "", meta


def _read_rtf(path: Path) -> Tuple[str, Dict[str, object]]:
    raw = path.read_text(encoding="latin-1", errors="ignore")
    text = _rtf_to_text(raw)
    meta: Dict[str, object] = {"rtf_characters": len(text)}
    return text, meta


def _read_eml(path: Path) -> Tuple[str, Dict[str, object]]:
    data = path.read_bytes()
    meta: Dict[str, object] = {}
    parser = BytesParser(policy=policy.default)
    try:
        message = parser.parsebytes(data)
    except Exception as exc:
        meta["email_parsed"] = False
        meta["email_error"] = str(exc)
        return "", meta

    meta.update(
        {
            "email_parsed": True,
            "email_subject": message.get("subject", ""),
            "email_from": message.get("from", ""),
            "email_to": message.get("to", ""),
            "email_date": message.get("date", ""),
        }
    )

    text_parts: List[str] = []
    html_fallback: List[str] = []
    attachments: List[Dict[str, object]] = []
    for part in message.walk():
        if part.is_multipart():
            continue
        filename = part.get_filename() or ""
        content_type = part.get_content_type()
        payload = part.get_payload(decode=True) or b""
        size = len(payload)
        if content_type.startswith("text/") and not filename:
            charset = part.get_content_charset() or "utf-8"
            try:
                decoded = payload.decode(charset, errors="ignore")
            except LookupError:
                decoded = payload.decode("utf-8", errors="ignore")
            if content_type == "text/plain":
                text_parts.append(decoded.strip())
            else:
                html_fallback.append(decoded)
        else:
            attachments.append(
                {
                    "filename": filename,
                    "content_type": content_type,
                    "size": size,
                }
            )

    if not text_parts and html_fallback:
        stripped = []
        for html in html_fallback:
            cleaned = re.sub(r"<[^>]+>", " ", html)
            cleaned = re.sub(r"\s+", " ", cleaned)
            stripped.append(cleaned.strip())
        text_parts = stripped

    meta["email_attachment_count"] = len(attachments)
    if attachments:
        meta["email_attachments"] = attachments[:10]
    meta["email_body_parts"] = len(text_parts)
    body = "\n\n".join(part for part in text_parts if part)
    return _normalize_newlines(body), meta


def _read_ics(path: Path) -> Tuple[str, Dict[str, object]]:
    raw, meta = _read_text_best_effort(path)
    text = _normalize_newlines(raw)
    events = text.upper().split("BEGIN:VEVENT")
    count = max(0, len(events) - 1)
    meta.update(
        {
            "ics_events": count,
            "ics_has_timezone": "TZID=" in text,
            "ics_has_alarms": "BEGIN:VALARM" in text,
        }
    )
    return text, meta


def _extract_image_text(path: Path) -> Tuple[str, Dict[str, object], List[Dict[str, Any]]]:
    """Return OCR/metadata text plus diagnostic details for *path*."""

    meta: Dict[str, object] = {}
    segments: List[Dict[str, Any]] = []
    text_pieces: List[str] = []
    text_sources: List[Dict[str, object]] = []
    ocr_segments = 0

    try:
        from PIL import ExifTags, Image  # type: ignore
    except Exception:
        meta["image_ocr_available"] = False
        return "", meta, segments

    try:
        with Image.open(str(path)) as im:
            corrected, ocr_ready, orientation_changed = _prepare_image_for_ocr(im)
            width, height = corrected.size
            meta.update(
                {
                    "image_mode": corrected.mode,
                    "image_size": (width, height),
                    "image_format": corrected.format or im.format or path.suffix.lstrip(".").upper(),
                }
            )
            if orientation_changed:
                meta["image_orientation_corrected"] = True

            info = getattr(im, "info", {}) or {}
            info_records: List[Dict[str, object]] = []
            for key, value in info.items():
                text_value: Optional[str] = None
                if isinstance(value, bytes):
                    try:
                        text_value = value.decode("utf-8", errors="ignore")
                    except Exception:
                        text_value = None
                elif isinstance(value, str):
                    text_value = value
                if not text_value:
                    continue
                cleaned = text_value.strip()
                if not cleaned:
                    continue
                info_records.append({"key": key, "length": len(cleaned)})
                text_pieces.append(cleaned)
                segments.append({"text": cleaned, "source": "metadata", "key": key, "kind": "metadata"})
            if info_records:
                text_sources.extend(
                    {"source": "metadata", "key": record["key"], "length": record["length"]}
                    for record in info_records[:10]
                )

            try:
                exif = im.getexif()
                if exif:
                    sample = {ExifTags.TAGS.get(k, str(k)): v for k, v in list(exif.items())[:40]}
                    if sample:
                        meta["image_exif_sample"] = sample
                    exif_records: List[Dict[str, object]] = []
                    for key, value in sample.items():
                        text_value: Optional[str] = None
                        if isinstance(value, bytes):
                            try:
                                text_value = value.decode("utf-8", errors="ignore")
                            except Exception:
                                text_value = None
                        elif isinstance(value, str):
                            text_value = value
                        if not text_value:
                            continue
                        cleaned = text_value.strip()
                        if not cleaned:
                            continue
                        exif_records.append({"key": key, "length": len(cleaned)})
                        text_pieces.append(cleaned)
                        segments.append({"text": cleaned, "source": "exif", "key": key, "kind": "metadata"})
                    if exif_records:
                        text_sources.extend(
                            {"source": "exif", "key": record["key"], "length": record["length"]}
                            for record in exif_records[:10]
                        )
            except Exception:
                pass

            try:
                import pytesseract  # type: ignore
                from pytesseract import Output  # type: ignore

                languages = _select_ocr_languages(pytesseract)
                lang_arg = "+".join(languages) if languages else ""
                if languages:
                    meta["image_ocr_languages"] = languages

                kwargs = {"output_type": Output.DICT}
                if lang_arg:
                    kwargs["lang"] = lang_arg

                data = pytesseract.image_to_data(ocr_ready, **kwargs)  # type: ignore[attr-defined]
                texts = data.get("text", [])
                if texts:
                    line_map: Dict[Tuple[int, int, int, int], Dict[str, Any]] = {}
                    total = len(texts)
                    for idx in range(total):
                        raw = texts[idx]
                        if not isinstance(raw, str):
                            continue
                        token = raw.strip()
                        if not token:
                            continue
                        page = int(data.get("page_num", [1])[idx]) - 1
                        block = int(data.get("block_num", [0])[idx])
                        paragraph = int(data.get("par_num", [0])[idx])
                        line = int(data.get("line_num", [0])[idx])
                        key = (page, block, paragraph, line)
                        entry = line_map.setdefault(
                            key,
                            {
                                "words": [],
                                "bbox": None,
                                "conf": [],
                                "page": page,
                            },
                        )
                        entry["words"].append(token)
                        left = float(data.get("left", [0])[idx])
                        top = float(data.get("top", [0])[idx])
                        width = float(data.get("width", [0])[idx])
                        height = float(data.get("height", [0])[idx])
                        bbox = entry["bbox"]
                        if bbox is None:
                            bbox = [left, top, left + width, top + height]
                        else:
                            bbox[0] = min(bbox[0], left)
                            bbox[1] = min(bbox[1], top)
                            bbox[2] = max(bbox[2], left + width)
                            bbox[3] = max(bbox[3], top + height)
                        entry["bbox"] = bbox
                        conf_values = entry["conf"]
                        conf_raw = data.get("conf", ["0"])[idx]
                        try:
                            conf_values.append(float(conf_raw))
                        except Exception:
                            pass

                    order = 1
                    for key, entry in sorted(line_map.items()):
                        words = entry["words"]
                        if not words:
                            continue
                        line_text = " ".join(words).strip()
                        if not line_text:
                            continue
                        bbox = entry["bbox"]
                        conf_values = entry["conf"]
                        avg_conf = 0.6
                        if conf_values:
                            avg_conf = sum(conf_values) / len(conf_values) / 100.0
                        page = entry.get("page", 0)
                        segments.append(
                            {
                                "text": line_text,
                                "source": "ocr",
                                "kind": "ocr",
                                "bbox": tuple(bbox) if bbox else None,
                                "confidence": round(avg_conf, 4),
                                "page": page,
                                "order": order,
                            }
                        )
                        text_pieces.append(line_text)
                        order += 1

                    if order > 1:
                        ocr_segments = order - 1
                        entry: Dict[str, Any] = {"source": "ocr", "lines": ocr_segments}
                        if languages:
                            entry["languages"] = languages
                        text_sources.append(entry)
                        meta["image_ocr_engine"] = "pytesseract"
            except ModuleNotFoundError:
                meta.setdefault("image_ocr_engine", "unavailable")
            except Exception as exc:
                meta["image_ocr_error"] = type(exc).__name__
    except Exception as exc:
        meta.setdefault("image_error", type(exc).__name__)

    if ocr_segments:
        meta["image_ocr_available"] = True
        meta["image_ocr_lines"] = ocr_segments
    else:
        meta.setdefault("image_ocr_available", False)

    if text_sources:
        meta["image_text_sources"] = text_sources[:10]

    deduped: List[str] = []
    seen: set[str] = set()
    for piece in text_pieces:
        if piece not in seen:
            deduped.append(piece)
            seen.add(piece)

    text = "\n".join(deduped).strip()
    if text:
        meta["image_has_text"] = True
        meta["image_text_characters"] = len(text)
        meta["image_text_preview"] = text[:400]
    else:
        meta.setdefault("image_has_text", False)

    return text, meta, segments


def _read_image(path: Path) -> Tuple[str, Dict[str, object]]:
    text, meta, _ = _extract_image_text(path)
    meta.update(_binary_preview(path.read_bytes()))
    return text, meta


def _read_zip_index(path: Path) -> Tuple[str, Dict[str, object]]:
    try:
        import zipfile
        with zipfile.ZipFile(str(path)) as zf:
            names = zf.namelist()
            meta = {"zip_entries": len(names), "zip_first_20": names[:20]}
            return "", meta
    except Exception as e:
        meta = _binary_preview(path.read_bytes())
        meta["zip_index"] = f"failed: {type(e).__name__}"
        return "", meta


# ------------------------- dispatcher -------------------------

def read_file_contents(path: Path, mime: str) -> Tuple[str, Dict[str, object]]:
    """Return a textual representation and extra metadata for *path*.

    Parameters
    ----------
    path: File to read.
    mime: MIME type determined by the caller (used as a hint).
    """
    _size_guard(path)
    suffix = path.suffix.lower()
    extra: Dict[str, object] = {}

    # JSON / JSONL
    if mime == "application/json" or suffix == ".json":
        return _read_json(path)
    if mime in _JSONL_MIMES or suffix in {".jsonl", ".ndjson"}:
        return _read_jsonl(path)
    if mime in {"application/x-yaml", "text/yaml", "application/yaml"} or suffix in {".yaml", ".yml"}:
        return _read_yaml(path)
    if mime in {"application/toml", "application/x-toml"} or suffix == ".toml":
        return _read_toml(path)
    if mime in {"application/xml", "text/xml"} or suffix == ".xml":
        return _read_xml(path)

    # HTML / Markdown / Plain text
    if mime in {"text/html"} or suffix in {".html", ".htm"}:
        return _read_html(path)
    if suffix in {".md", ".markdown"}:
        return _read_markdown(path)
    if mime in {"text/rtf"} or suffix == ".rtf":
        return _read_rtf(path)
    if mime == "message/rfc822" or suffix == ".eml":
        return _read_eml(path)
    if mime == "text/calendar" or suffix == ".ics":
        return _read_ics(path)
    if mime.startswith("text/") or suffix in _TEXT_EXTENSIONS:
        return _read_plain_text(path)

    # PDF / Office
    if mime == "application/pdf" or suffix == ".pdf":
        return _read_pdf(path)
    if suffix == ".docx":
        return _read_docx(path)
    if suffix == ".pptx":
        return _read_pptx(path)
    if suffix == ".xlsx":
        return _read_xlsx(path)

    # Images
    if mime.startswith("image/"):
        return _read_image(path)

    # Archives
    if suffix == ".zip":
        return _read_zip_index(path)

    # Fallback: binary preview only
    blob = path.read_bytes()
    extra.update(_binary_preview(blob))
    return "", extra
